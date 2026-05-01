from pathlib import Path
from io import BytesIO
from typing import List, Dict, Any
try:
    import pymupdf as fitz
except Exception:
    import fitz
import easyocr
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import docx2txt
from PIL import Image
import numpy as np
from src.config import SUPPORTED_EXTENSIONS


_ocr_reader = None


def get_ocr_reader():
    global _ocr_reader
    if _ocr_reader is None:
        print("Loading EasyOCR model (first time may download ~100MB)...")
        _ocr_reader = easyocr.Reader(["en"], gpu=False)
        print("EasyOCR ready")
    return _ocr_reader


def ocr_image_bytes(image_bytes: bytes) -> str:
    reader = get_ocr_reader()
    image = Image.open(BytesIO(image_bytes))
    img_array = np.array(image)
    results = reader.readtext(img_array, detail=0)
    return "\n".join(results)


class DocumentLoader:

    def __init__(self, data_dir: str = "."):
        self.data_dir = Path(data_dir).resolve()

    def _load_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        elements = []
        if not hasattr(fitz, "open"):
            raise RuntimeError(
                "PDF engine is not available. Install PyMuPDF (`pip install pymupdf`) "
                "and remove any conflicting `fitz` package."
            )
        pdf = fitz.open(file_path)

        for page_num, page in enumerate(pdf, start=1):
            text = page.get_text()
            if text.strip():
                elements.append({
                    "type": "text",
                    "content": text.strip(),
                    "metadata": {"source": file_path, "page": page_num},
                })

            for img_index, img_info in enumerate(page.get_images(full=True)):
                try:
                    xref = img_info[0]
                    base_image = pdf.extract_image(xref)
                    image_bytes = base_image["image"]
                    ocr_text = ocr_image_bytes(image_bytes)
                    if ocr_text.strip():
                        elements.append({
                            "type": "image",
                            "content": ocr_text.strip(),
                            "metadata": {
                                "source": file_path,
                                "page": page_num,
                                "image_index": img_index,
                            },
                        })
                except Exception:
                    pass

        pdf.close()
        return elements

    def _load_docx(self, file_path: str) -> List[Dict[str, Any]]:
        text = docx2txt.process(file_path)
        if text.strip():
            return [{
                "type": "text",
                "content": text.strip(),
                "metadata": {"source": file_path},
            }]
        return []

    def _load_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        elements = []
        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            tables = []
            image_texts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)

                if shape.has_table:
                    table = shape.table
                    header = " | ".join(cell.text for cell in table.rows[0].cells)
                    sep = " | ".join("---" for _ in table.rows[0].cells)
                    rows = [header, sep]
                    for row in list(table.rows)[1:]:
                        rows.append(" | ".join(cell.text for cell in row.cells))
                    table_md = "\n".join(rows)
                    tables.append(table_md)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        ocr_text = ocr_image_bytes(shape.image.blob)
                        if ocr_text.strip():
                            image_texts.append(ocr_text.strip())
                    except Exception:
                        pass

                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub in shape.shapes:
                        if hasattr(sub, "text_frame") and sub.has_text_frame:
                            texts.append(sub.text_frame.text)
                        if hasattr(sub, "image"):
                            try:
                                ocr = ocr_image_bytes(sub.image.blob)
                                if ocr.strip():
                                    image_texts.append(ocr.strip())
                            except Exception:
                                pass

            if texts:
                elements.append({
                    "type": "text",
                    "content": "\n".join(texts),
                    "metadata": {"source": file_path, "slide": slide_num},
                })
            for t in tables:
                elements.append({
                    "type": "table",
                    "content": t,
                    "metadata": {"source": file_path, "slide": slide_num},
                })
            for it in image_texts:
                elements.append({
                    "type": "image",
                    "content": it,
                    "metadata": {"source": file_path, "slide": slide_num},
                })

        return elements

    def _load_file(self, file_path: Path) -> List[Dict[str, Any]]:
        ext = file_path.suffix.lower()
        path_str = str(file_path)

        if ext == ".pdf":
            return self._load_pdf(path_str)
        elif ext in {".docx", ".doc"}:
            return self._load_docx(path_str)
        elif ext in {".pptx", ".ppt"}:
            return self._load_pptx(path_str)
        else:
            raise ValueError(f"Unsupported: {ext}")

    def load_all(self) -> List[Dict[str, Any]]:
        all_elements = []
        files = [
            f for f in self.data_dir.rglob("*")
            if f.suffix.lower() in SUPPORTED_EXTENSIONS
        ]

        print(f"Found {len(files)} supported file(s)")
        for f in files:
            print(f"  {f.name}")

        for file in files:
            try:
                elements = self._load_file(file)
                all_elements.extend(elements)
                print(f"  Loaded: {file.name} -> {len(elements)} elements")
            except Exception as e:
                print(f"  Failed: {file.name} -> {e}")

        print(f"Total elements: {len(all_elements)}")
        return all_elements

    def load_single(self, file_path: str) -> List[Dict[str, Any]]:
        fp = Path(file_path).resolve()
        if fp.suffix.lower() not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported: {fp.suffix}")

        print(f"Loading: {fp.name}")
        elements = self._load_file(fp)
        print(f"{len(elements)} elements extracted")
        return elements